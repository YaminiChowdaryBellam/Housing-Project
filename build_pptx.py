"""
Build the HACR Housing Analysis PowerPoint presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = "/Users/yamini/Documents/Housing-Project/Wake_County_Housing_Analysis.pptx"
CHARTS = "/Users/yamini/Documents/Housing-Project/charts"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
TEAL   = RGBColor(0x00, 0x7A, 0x87)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
AMBER  = RGBColor(0xE6, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
LIGHT  = RGBColor(0xE8, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines):
    """lines = list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def kpi_card(slide, left, top, number, label, sub, card_color, text_color=WHITE):
    """Draw a KPI card with big number, label, and subtitle."""
    card_w = Inches(2.8)
    card_h = Inches(2.5)
    shape = add_shape(slide, left, top, card_w, card_h, card_color)

    # Big number
    add_text_box(slide, left + Inches(0.2), top + Inches(0.25),
                 card_w - Inches(0.4), Inches(1.0),
                 number, font_size=36, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.2), top + Inches(1.2),
                 card_w - Inches(0.4), Inches(0.6),
                 label, font_size=13, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    add_text_box(slide, left + Inches(0.15), top + Inches(1.75),
                 card_w - Inches(0.3), Inches(0.65),
                 sub, font_size=9, color=text_color, bold=False, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Accent bar top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

# Title
add_multi_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(3.5), [
    ("Closing the Gap", 44, WHITE, True, PP_ALIGN.CENTER),
    ("Wake County Affordable Housing Supply-Demand Analysis", 22, TEAL, False, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("Built with Real Data from wakehousingdata.org", 16, WHITE, False, PP_ALIGN.CENTER),
    ("Wake County HACR  |  April 2026", 14, RGBColor(0xAA, 0xBB, 0xCC), False, PP_ALIGN.CENTER),
])

# Bottom accent
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — KPI Dashboard
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "KEY PERFORMANCE INDICATORS  |  Wake County Housing Crisis at a Glance",
             font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

# Row 1: 4 KPI cards
kpi_card(slide, Inches(0.5), Inches(1.3),
         "26,037", "Units Missing", "Affordable rental shortage\nfor households under $35K",
         RED)

kpi_card(slide, Inches(3.65), Inches(1.3),
         "39 Years", "To Close the Gap", "At current HACR pace of\n660 units/year (closes 2065)",
         NAVY)

kpi_card(slide, Inches(6.8), Inches(1.3),
         "$3.8B", "Total Investment Needed", "Same cost regardless of pace\u2014\nonly the timeline changes",
         AMBER)

kpi_card(slide, Inches(9.95), Inches(1.3),
         "89%", "Entry-Level Collapse", "Drop in homes sold under $250K\n2018 to 2022",
         RGBColor(0x5B, 0x2D, 0x8E))

# Row 2: 4 more KPI cards
kpi_card(slide, Inches(0.5), Inches(4.2),
         "3.72%", "Vacancy Rate", "Well below 5% healthy benchmark\u2014\nlandlords set the terms",
         TEAL)

kpi_card(slide, Inches(3.65), Inches(4.2),
         "+$8,920", "2021 Rent Shock", "Income jump needed in a single\nyear to keep up with rent",
         RED)

kpi_card(slide, Inches(6.8), Inches(4.2),
         "96%", "Renters Under $20K", "Are cost-burdened\u2014spending\n>30% of income on housing",
         RGBColor(0x8B, 0x0A, 0x0A))

kpi_card(slide, Inches(9.95), Inches(4.2),
         "16 / 100", "Affordable Homes", "Only 16 affordable units for\nevery 100 renters under $20K",
         DARK)

# Footer
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             "Source: wakehousingdata.org  |  HACR FY2025 Annual Report  |  All real data, no estimates",
             font_size=9, color=RGBColor(0x88, 0x99, 0xAA), bold=False, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Gap
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "THE GAP  |  Rental Housing Deficit by Income Band",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/01_supply_demand_gap.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

# Insight box right side
add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), LIGHT, TEAL)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("What This Shows", 16, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("Red bars = shortage. Green bars = surplus.", 12, DARK, False, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("The shortage is concentrated at the lowest income levels:", 11, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("  <$20K:  16,774 units short", 11, RED, True, PP_ALIGN.LEFT),
    ("  <$35K:  26,037 units short", 11, RED, True, PP_ALIGN.LEFT),
    ("  <$50K:  24,468 units short", 11, RED, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("Above $75K, supply exceeds demand. The crisis is an income-targeted problem.", 11, DARK, False, PP_ALIGN.LEFT),
])

# Bottom insight strip
add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  Wake County has enough housing overall \u2014 but not for people who need it most. "
             "The market works for $75K+ households. Below $50K, 67,279 families have no affordable option.",
             font_size=13, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Affordable Units per 100
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), TEAL)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "SUPPLY PER 100 RENTERS  |  How Many Chairs for 100 People?",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/02_per_100_renters.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.0))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.0), LIGHT, TEAL)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.6), [
    ("Musical Chairs Analogy", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("At <$20K income:", 12, DARK, True, PP_ALIGN.LEFT),
    ("100 people, only 16 chairs.", 12, RED, True, PP_ALIGN.LEFT),
    ("84 people are left standing.", 12, RED, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("At <$75K income:", 12, DARK, True, PP_ALIGN.LEFT),
    ("100 people, 104 chairs.", 12, GREEN, True, PP_ALIGN.LEFT),
    ("Everyone sits, a few to spare.", 12, GREEN, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("The dashed line = 100 (balanced).", 11, DARK, False, PP_ALIGN.LEFT),
    ("Below it = shortage. Above it = surplus.", 11, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.0), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.6), Inches(12), Inches(0.8),
             "KEY INSIGHT:  The lowest-income renters face the most extreme competition for housing. "
             "At <$20K, each affordable unit has 6 households competing for it.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Cost Burden Explosion
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), RED)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "COST BURDEN EXPLOSION  |  The Crisis Moves Up the Income Ladder",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/03_cost_burden.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), RGBColor(0xFD, 0xE8, 0xE8), RED)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("Biggest Shifts (2014 to 2024)", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("$35K-$50K:  34.5%  -->  88.9%", 12, RED, True, PP_ALIGN.LEFT),
    ("Jumped +54 percentage points", 11, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("$50K-$75K:  9.8%  -->  55.8%", 12, RED, True, PP_ALIGN.LEFT),
    ("Nearly 6x increase", 11, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("$75K-$100K:  2.7%  -->  20.8%", 12, AMBER, True, PP_ALIGN.LEFT),
    ("Even upper-middle income is affected", 11, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("The dashed line = HUD's 30% threshold.", 10, DARK, False, PP_ALIGN.LEFT),
    ("Every bar above it = cost-burdened.", 10, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  In 2014, only renters under $35K were in crisis. By 2024, the crisis has engulfed "
             "the $35K-$75K range \u2014 the working class and middle class are now cost-burdened too.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Income vs Rent Gap
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "THE AFFORDABILITY SQUEEZE  |  Renter Income vs. What Rent Actually Costs",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/04_income_vs_rent.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), LIGHT, TEAL)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("Reading This Chart", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("Teal line = what renters actually earn.", 11, TEAL, False, PP_ALIGN.LEFT),
    ("Red line = income needed to afford rent.", 11, RED, False, PP_ALIGN.LEFT),
    ("Shaded area = the gap.", 11, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("The 2021 Shock:", 12, RED, True, PP_ALIGN.LEFT),
    ("Rents jumped $223/month in one year.", 11, DARK, False, PP_ALIGN.LEFT),
    ("Income needed to jump $8,920.", 11, DARK, False, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("By 2024:", 12, GREEN, True, PP_ALIGN.LEFT),
    ("The lines nearly converge \u2014 the gap is", 11, DARK, False, PP_ALIGN.LEFT),
    ("closing, but only because incomes are", 11, DARK, False, PP_ALIGN.LEFT),
    ("finally catching up after 3 years.", 11, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  The median renter in Wake County spent 3 years underwater after the 2021 rent shock. "
             "The gap peaked at $8,272/year \u2014 money that came from food, healthcare, and savings.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Entry-Level Collapse
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), RGBColor(0x5B, 0x2D, 0x8E))
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "ENTRY-LEVEL COLLAPSE  |  First-Time Buyers Priced Out",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/05_entry_level_collapse.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), RGBColor(0xF0, 0xEA, 0xF8), RGBColor(0x5B, 0x2D, 0x8E))
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("The Shift in 4 Years", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("Under $250K homes:", 12, DARK, True, PP_ALIGN.LEFT),
    ("2018: 28.6% of all sales", 12, TEAL, True, PP_ALIGN.LEFT),
    ("2022: 3.1% of all sales", 12, RED, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("$500K+ homes:", 12, DARK, True, PP_ALIGN.LEFT),
    ("2018: 17.7% of all sales", 12, TEAL, False, PP_ALIGN.LEFT),
    ("2022: 46.8% of all sales", 12, RGBColor(0x5B, 0x2D, 0x8E), True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("The market flipped: nearly half of all", 11, DARK, False, PP_ALIGN.LEFT),
    ("homes sold are now $500K+, a price", 11, DARK, False, PP_ALIGN.LEFT),
    ("most first-time buyers cannot reach.", 11, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  In 2018, nearly 1 in 3 homes sold was affordable for a first-time buyer. "
             "By 2022, it was 1 in 32. The path from renting to owning has effectively closed for most Wake County residents.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Household Displacement
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "DISPLACEMENT  |  Who Left Wake County? (2014-2024)",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/07_hh_change.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), LIGHT, TEAL)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("The Wealth Polarisation", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("Bottom 3 bands lost:", 12, DARK, True, PP_ALIGN.LEFT),
    ("-30,353 households", 14, RED, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("Top 2 bands gained:", 12, DARK, True, PP_ALIGN.LEFT),
    ("+129,369 households", 14, GREEN, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("4:1 ratio \u2014 for every low-income", 11, DARK, False, PP_ALIGN.LEFT),
    ("household that left, four high-income", 11, DARK, False, PP_ALIGN.LEFT),
    ("households moved in.", 11, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  23,176 households earning under $35K disappeared from Wake County in a decade. "
             "They didn't get richer \u2014 they were priced out. Meanwhile, $150K+ households grew by 100,784.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Rent Trend + Vacancy
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), TEAL)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "SUPPLY PRESSURE  |  Rent Trend & Vacancy Rate",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/09_rent_trend.png",
                         Inches(0.3), Inches(1.1), Inches(6.3), Inches(3.5))

slide.shapes.add_picture(f"{CHARTS}/06_vacancy_rate.png",
                         Inches(6.8), Inches(1.1), Inches(6.3), Inches(3.5))

add_shape(slide, Inches(0.4), Inches(4.9), Inches(12.5), Inches(1.8), NAVY)
add_multi_text(slide, Inches(0.7), Inches(5.0), Inches(12), Inches(1.6), [
    ("How These Two Charts Connect:", 14, AMBER, True, PP_ALIGN.LEFT),
    ("Low vacancy (left chart: 2.74% in 2019) gave landlords pricing power \u2192 rents spiked "
     "+$223/month in 2021 (right chart). As new multifamily units came online (2022+), vacancy "
     "rose to 3.77% and rent growth stalled. But 3.72% is still well below the 5% healthy "
     "benchmark \u2014 tenants have no real bargaining power.", 12, WHITE, False, PP_ALIGN.LEFT),
])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Homeownership by Race
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "EQUITY LENS  |  Homeownership Rate by Race / Ethnicity",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/10_ownership_by_race.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), LIGHT, TEAL)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("The Equity Gap", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("White homeownership:  71.4%", 12, NAVY, True, PP_ALIGN.LEFT),
    ("Black homeownership:  45.5%", 12, RED, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("25.9 percentage point gap.", 13, RED, True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("With 88,430 Black households in", 11, DARK, False, PP_ALIGN.LEFT),
    ("Wake County, that gap represents", 11, DARK, False, PP_ALIGN.LEFT),
    ("roughly 22,900 families locked out", 11, DARK, False, PP_ALIGN.LEFT),
    ("of homeownership relative to the", 11, DARK, False, PP_ALIGN.LEFT),
    ("White rate.", 11, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  Wake County's homeownership gap mirrors national disparities. "
             "Black and Hispanic/Latino households own at rates 20-26 points below White households \u2014 "
             "a gap that compounds generational wealth inequality.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Scenario Planner
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), AMBER)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "SCENARIO PLANNER  |  What Would It Take to Close the Gap?",
             font_size=20, color=WHITE, bold=True)

slide.shapes.add_picture(f"{CHARTS}/08_scenarios.png",
                         Inches(0.4), Inches(1.1), Inches(7.5), Inches(4.2))

add_shape(slide, Inches(8.3), Inches(1.1), Inches(4.6), Inches(4.2), RGBColor(0xFF, 0xF8, 0xE8), AMBER)
add_multi_text(slide, Inches(8.5), Inches(1.3), Inches(4.2), Inches(3.8), [
    ("4 Policy Scenarios", 15, NAVY, True, PP_ALIGN.LEFT),
    ("", 8, DARK, False, PP_ALIGN.LEFT),
    ("Current Pace (660/yr):", 11, DARK, True, PP_ALIGN.LEFT),
    ("40 years. Closes in 2065.", 11, RED, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("2x Scale-Up (1,320/yr):", 11, DARK, True, PP_ALIGN.LEFT),
    ("20 years. Closes in 2045.", 11, AMBER, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("5x Scale-Up (3,300/yr):", 11, DARK, True, PP_ALIGN.LEFT),
    ("8 years. Closes in 2033.", 11, GREEN, True, PP_ALIGN.LEFT),
    ("", 4, DARK, False, PP_ALIGN.LEFT),
    ("Deep Subsidy (<$20K focus):", 11, DARK, True, PP_ALIGN.LEFT),
    ("40 years. Prioritizes deepest need.", 11, RGBColor(0x5B, 0x2D, 0x8E), True, PP_ALIGN.LEFT),
    ("", 6, DARK, False, PP_ALIGN.LEFT),
    ("Total cost: $3.8B in ALL scenarios.", 12, AMBER, True, PP_ALIGN.LEFT),
    ("Only the annual budget changes.", 11, DARK, False, PP_ALIGN.LEFT),
])

add_shape(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.1), NAVY)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(12), Inches(0.9),
             "KEY INSIGHT:  The total investment is fixed at $3.8B. The only question is how fast to spend it. "
             "Doubling production halves the timeline \u2014 but the county will spend the same total amount regardless.",
             font_size=13, color=WHITE, bold=False)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Data Model
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             "DATA MODEL  |  How the 4 Source Files Connect",
             font_size=20, color=WHITE, bold=True)

# Add the data model image
dm_path = "/Users/yamini/Documents/Housing-Project/data_model_diagram.png"
if os.path.exists(dm_path):
    slide.shapes.add_picture(dm_path, Inches(0.8), Inches(1.1), Inches(11.7), Inches(5.8))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Summary & Next Steps
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             "SUMMARY & NEXT STEPS", font_size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Left column — What We Found
add_shape(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.0), RGBColor(0x24, 0x4D, 0x82))
add_multi_text(slide, Inches(0.9), Inches(1.7), Inches(5.3), Inches(4.6), [
    ("What the Data Tells Us", 16, AMBER, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("1.  26,037 affordable units missing for <$35K renters", 12, WHITE, False, PP_ALIGN.LEFT),
    ("2.  39 years to close at current pace (660 units/yr)", 12, WHITE, False, PP_ALIGN.LEFT),
    ("3.  $3.8B total investment needed regardless of timeline", 12, WHITE, False, PP_ALIGN.LEFT),
    ("4.  Entry-level homes collapsed 89% (2018-2022)", 12, WHITE, False, PP_ALIGN.LEFT),
    ("5.  Crisis has spread to $35K-$75K income bands", 12, WHITE, False, PP_ALIGN.LEFT),
    ("6.  30,353 low-income HHs displaced in a decade", 12, WHITE, False, PP_ALIGN.LEFT),
    ("7.  Vacancy at 3.72% gives renters zero leverage", 12, WHITE, False, PP_ALIGN.LEFT),
    ("8.  25.9 pp racial homeownership gap persists", 12, WHITE, False, PP_ALIGN.LEFT),
])

# Right column — Recommendations
add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.0), RGBColor(0x24, 0x4D, 0x82))
add_multi_text(slide, Inches(7.2), Inches(1.7), Inches(5.3), Inches(4.6), [
    ("What I Would Build Next", 16, TEAL, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("1.  Census tract-level mapping (geographic equity)", 12, WHITE, False, PP_ALIGN.LEFT),
    ("2.  Race/ethnicity overlay on all gap analyses", 12, WHITE, False, PP_ALIGN.LEFT),
    ("3.  Live HMIS integration (homelessness correlation)", 12, WHITE, False, PP_ALIGN.LEFT),
    ("4.  \"What-if\" tool for new development proposals", 12, WHITE, False, PP_ALIGN.LEFT),
    ("5.  Automated CAPER table generation from live data", 12, WHITE, False, PP_ALIGN.LEFT),
    ("6.  Landlord incentive ROI calculator for HACR", 12, WHITE, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Built with:", 12, AMBER, True, PP_ALIGN.LEFT),
    ("Python 3.9  |  Streamlit  |  Plotly  |  pandas", 12, WHITE, False, PP_ALIGN.LEFT),
    ("All data from wakehousingdata.org (real, public)", 12, WHITE, False, PP_ALIGN.LEFT),
])

# Bottom bar
add_shape(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), AMBER)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"\u2713 Saved: {OUTPUT}")
print(f"  13 slides | 10 chart visuals | KPI dashboard | Data model")
